"""Tests for the extractor protocol & built-in extractors (RFC 004 §6)."""

from __future__ import annotations

import io
from pathlib import Path

import pytest

from gyra.knowledge.types import ExtractMode
from gyra_ext.knowledge.extractors import (
    Extractor,
    VerbatimSpec,
    get_extractor_registry,
)
from gyra_ext.knowledge.extractors import builtin as builtin_mod
from gyra_ext.knowledge.extractors.registry_init import register_builtin_extractors


@pytest.fixture(autouse=True, scope="module")
def _ensure_builtins_registered():
    """Make sure built-ins are registered once for the whole module.

    The `@extractor` decorator runs at `builtin.py` import time and is
    idempotent (re-registrations overwrite by name), so we don't need to
    clear the registry between tests.
    """
    register_builtin_extractors()
    yield


def test_registry_has_builtins():
    reg = get_extractor_registry()
    names = {e.name for e in reg.list_all()}
    assert {
        "text", "pdf", "docx", "pptx", "excel", "image", "audio", "video",
    }.issubset(names)


def test_registry_get_by_mime():
    reg = get_extractor_registry()
    assert reg.get("text/plain").name == "text"
    assert reg.get("text/markdown").name == "text"
    assert reg.get("application/json").name == "text"
    assert reg.get("application/pdf").name == "pdf"
    assert reg.get("image/png").name == "image"
    assert reg.get("audio/mpeg").name == "audio"
    assert reg.get("application/vnd.ms-excel").name == "excel"
    xlsx_mime = (
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    )
    assert reg.get(xlsx_mime).name == "excel"
    assert reg.get("application/octet-stream") is None


@pytest.mark.asyncio
async def test_text_extractor_basic(tmp_path: Path):
    reg = get_extractor_registry()
    ext = reg.get("text/plain")
    assert isinstance(ext, Extractor)

    f = tmp_path / "note.txt"
    f.write_text("hello world\nsecond line", encoding="utf-8")

    specs = await ext.extract(f, "text/plain", model=None, model_caller=None)
    assert len(specs) == 1
    spec = specs[0]
    assert isinstance(spec, VerbatimSpec)
    assert spec.content == "hello world\nsecond line"
    assert spec.source_file == "note.txt"
    assert spec.extract_mode == ExtractMode.UPLOAD
    assert spec.meta.get("mime") == "text/plain"
    assert spec.meta.get("bytes") == len(f.read_bytes())
    assert spec.content_date is not None


@pytest.mark.asyncio
async def test_text_extractor_handles_non_utf8(tmp_path: Path):
    """Non-UTF-8 bytes should not raise — fallback to replace."""
    ext = get_extractor_registry().get("text/plain")
    f = tmp_path / "bad.bin"
    f.write_bytes(b"\xff\xfe\x00bad bytes here")
    specs = await ext.extract(f, "text/plain", model=None, model_caller=None)
    assert len(specs) == 1
    assert "bad bytes here" in specs[0].content


@pytest.mark.asyncio
async def test_image_extractor_invokes_model_caller(tmp_path: Path, monkeypatch):
    """ImageExtractor should hand off to the supplied model_caller, not decode bytes."""
    ext = get_extractor_registry().get("image/png")

    # Minimal fake PNG (8-byte signature + IHDR chunk). We don't need a real
    # image — ImageExtractor just reads bytes and passes the path to the caller.
    png_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 32
    f = tmp_path / "snap.png"
    f.write_bytes(png_bytes)

    captured: dict = {}

    async def fake_caller(model, prompt, images=None):
        captured["model"] = model
        captured["prompt"] = prompt
        captured["images"] = images or []
        return "A photo of a cat sitting on a desk."

    specs = await ext.extract(
        f, "image/png", model="gpt-4o", model_caller=fake_caller
    )
    assert len(specs) == 1
    spec = specs[0]
    assert spec.extract_mode == ExtractMode.UPLOAD
    assert "cat" in spec.content
    assert spec.source_file == "snap.png"
    assert captured["model"] == "gpt-4o"
    assert captured["images"] == [f]


@pytest.mark.asyncio
async def test_image_extractor_without_model_caller_writes_empty_caption(tmp_path: Path):
    """Without a model_caller, ImageExtractor falls back to an empty-caption verbatim."""
    ext = get_extractor_registry().get("image/png")
    f = tmp_path / "snap.png"
    f.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 16)

    specs = await ext.extract(f, "image/png", model=None, model_caller=None)
    assert len(specs) == 1
    # Caption line is empty
    assert "Caption" in specs[0].content
    assert specs[0].meta.get("caption_len") == 0


@pytest.mark.asyncio
async def test_pdf_extractor_missing_dep_raises_clear_error(tmp_path: Path, monkeypatch):
    """If pdfplumber isn't installed, the extractor should raise a clear message."""
    ext = get_extractor_registry().get("application/pdf")
    f = tmp_path / "doc.pdf"
    f.write_bytes(b"%PDF-1.4 fake pdf")

    # Force the import to fail
    import sys

    monkeypatch.setitem(sys.modules, "pdfplumber", None)

    with pytest.raises(RuntimeError, match="pdfplumber"):
        await ext.extract(f, "application/pdf", model=None, model_caller=None)


def test_register_builtin_is_idempotent():
    """Calling register_builtin_extractors twice should not duplicate entries."""
    reg = get_extractor_registry()
    count_before = len(reg.list_all())
    register_builtin_extractors()
    count_after = len(reg.list_all())
    assert count_before == count_after


# ---------------------------------------------------------------------------
# VideoExtractor: native direct + frame-pipeline fallback
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_video_extractor_native_strategy(tmp_path: Path):
    """Native video models (gemini/qwen-vl) receive the file via videos kwarg."""
    ext = get_extractor_registry().get("video/mp4")
    assert ext is not None and ext.name == "video"

    f = tmp_path / "clip.mp4"
    f.write_bytes(b"\x00" * 1024)  # small fake file, content never read by fake caller

    captured: dict = {}

    async def fake_caller(model, prompt, images=None, videos=None):
        captured["videos"] = videos or []
        captured["images"] = images or []
        return "一段产品介绍视频,演示了上传流程。"

    specs = await ext.extract(f, "video/mp4", model="gemini-2.5-pro", model_caller=fake_caller)
    assert len(specs) == 1
    assert specs[0].meta["strategy"] == "native"
    assert "上传流程" in specs[0].content
    assert captured["videos"] == [f]
    assert captured["images"] == []


@pytest.mark.asyncio
async def test_video_extractor_falls_back_to_frames(tmp_path: Path, monkeypatch):
    """Non-native model (or native failure) → ffmpeg frame pipeline via image channel."""
    from gyra_ext.knowledge.extractors.builtin import VideoExtractor

    ext = get_extractor_registry().get("video/mp4")
    f = tmp_path / "clip.mp4"
    f.write_bytes(b"\x00" * 1024)

    # 假帧文件:管线只把路径交给 model_caller,不读内容
    frames = []
    for i in range(3):
        fp = tmp_path / f"frame_{i:03d}.jpg"
        fp.write_bytes(b"\xff\xd8\xff" + b"\x00" * 16)
        frames.append(fp)

    async def fake_frames(self, path, out_dir):
        return frames

    async def fake_audio(self, path, out_dir):
        return None

    monkeypatch.setattr(VideoExtractor, "_extract_frames", fake_frames)
    monkeypatch.setattr(VideoExtractor, "_extract_audio", fake_audio)
    monkeypatch.setattr(builtin_mod.shutil, "which", lambda name: "/usr/bin/ffmpeg")

    captured: dict = {}

    async def fake_caller(model, prompt, images=None, videos=None):
        captured["images"] = images or []
        captured["videos"] = videos or []
        return "帧序列描述了会议投屏内容。"

    # gpt-4o 不命中原生模式 → 直接走抽帧
    specs = await ext.extract(f, "video/mp4", model="gpt-4o", model_caller=fake_caller)
    assert len(specs) == 1
    assert specs[0].meta["strategy"] == "frames"
    assert specs[0].meta["frames"] == 3
    assert captured["images"] == frames
    assert captured["videos"] == []


@pytest.mark.asyncio
async def test_video_extractor_native_failure_falls_back(tmp_path: Path, monkeypatch):
    """Native-capable model that errors must degrade to the frame pipeline."""
    from gyra_ext.knowledge.extractors.builtin import VideoExtractor

    ext = get_extractor_registry().get("video/mp4")
    f = tmp_path / "clip.mp4"
    f.write_bytes(b"\x00" * 1024)

    frame = tmp_path / "frame_001.jpg"
    frame.write_bytes(b"\xff\xd8\xff" + b"\x00" * 16)

    async def fake_frames(self, path, out_dir):
        return [frame]

    async def fake_audio(self, path, out_dir):
        return None

    monkeypatch.setattr(VideoExtractor, "_extract_frames", fake_frames)
    monkeypatch.setattr(VideoExtractor, "_extract_audio", fake_audio)
    monkeypatch.setattr(builtin_mod.shutil, "which", lambda name: "/usr/bin/ffmpeg")

    calls: list = []

    async def flaky_caller(model, prompt, images=None, videos=None):
        calls.append({"images": images or [], "videos": videos or []})
        if videos:
            raise RuntimeError("provider rejected video_url")
        return "抽帧后的画面描述。"

    specs = await ext.extract(f, "video/mp4", model="qwen2.5-vl-72b", model_caller=flaky_caller)
    assert len(specs) == 1
    assert specs[0].meta["strategy"] == "frames"
    # 先尝试 native(videos),失败后走 images
    assert calls[0]["videos"] == [f]
    assert calls[1]["images"] == [frame]


@pytest.mark.asyncio
async def test_video_extractor_requires_model(tmp_path: Path):
    """No model/model_caller → clear error (unlike image, video has no text fallback)."""
    ext = get_extractor_registry().get("video/mp4")
    f = tmp_path / "clip.mp4"
    f.write_bytes(b"\x00" * 64)

    with pytest.raises(RuntimeError, match="multimodal model"):
        await ext.extract(f, "video/mp4", model=None, model_caller=None)


# ---------------------------------------------------------------------------
# Embedded images in docx/pptx (asset_store persistence + vision caption)
# ---------------------------------------------------------------------------

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def _noise_png(width: int = 96, height: int = 96) -> bytes:
    """A real PNG (python-docx/pptx must parse it) large enough for the extractor."""
    import io
    import random

    from PIL import Image

    img = Image.new("RGB", (width, height))
    rng = random.Random(42)
    img.putdata(
        [
            (rng.randrange(256), rng.randrange(256), rng.randrange(256))
            for _ in range(width * height)
        ]
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    data = buf.getvalue()
    assert len(data) >= builtin_mod.MIN_EMBEDDED_IMAGE_BYTES
    return data


def _tiny_png() -> bytes:
    """A real but tiny PNG that must be skipped by MIN_EMBEDDED_IMAGE_BYTES."""
    import io

    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (1, 1)).save(buf, format="PNG")
    return buf.getvalue()


@pytest.mark.asyncio
async def test_docx_extractor_persists_and_captions_images(tmp_path: Path):
    pytest.importorskip("docx")
    pytest.importorskip("PIL")
    import docx as docx_mod

    png = _noise_png()
    doc = docx_mod.Document()
    doc.add_paragraph("第一段正文")
    doc.add_picture(io.BytesIO(png))
    doc.add_paragraph("第二段正文")
    f = tmp_path / "with_img.docx"
    doc.save(str(f))

    stored: dict = {}

    async def fake_asset_store(filename: str, data: bytes) -> str:
        stored[filename] = data
        return f"assets/fake-{filename}"

    captured: dict = {}

    async def fake_caller(model, prompt, images=None):
        captured["model"] = model
        captured["images"] = images or []
        # tmp file only exists during the call — read it now
        captured["captioned_bytes"] = images[0].read_bytes() if images else None
        return "一张测试图"

    ext = get_extractor_registry().get(DOCX_MIME)
    specs = await ext.extract(
        f,
        DOCX_MIME,
        model="vl-model",
        model_caller=fake_caller,
        asset_store=fake_asset_store,
    )

    assert len(specs) == 1
    spec = specs[0]
    assert spec.meta["images"] == 1
    assert "第一段正文" in spec.content
    assert "第二段正文" in spec.content
    assert "![图片1: image1.png](assets/fake-image1.png)" in spec.content
    assert "图片说明：一张测试图" in spec.content
    assert stored == {"image1.png": png}
    assert captured["model"] == "vl-model"
    assert captured["captioned_bytes"] == png


@pytest.mark.asyncio
async def test_docx_extractor_without_model_and_store_keeps_placeholder(tmp_path: Path):
    """No asset_store + no model → bare placeholder, document still ingests."""
    pytest.importorskip("docx")
    pytest.importorskip("PIL")
    import docx as docx_mod

    doc = docx_mod.Document()
    doc.add_picture(__import__("io").BytesIO(_noise_png()))
    f = tmp_path / "bare.docx"
    doc.save(str(f))

    ext = get_extractor_registry().get(DOCX_MIME)
    specs = await ext.extract(f, DOCX_MIME, model=None, model_caller=None)

    assert len(specs) == 1
    content = specs[0].content
    assert "（内嵌图片: image1.png，未保存）" in content
    assert "![图片1" not in content
    assert "图片说明：" not in content
    assert specs[0].meta["images"] == 1


@pytest.mark.asyncio
async def test_docx_extractor_skips_tiny_images(tmp_path: Path):
    """Images below MIN_EMBEDDED_IMAGE_BYTES are decoration → skipped entirely."""
    pytest.importorskip("docx")
    pytest.importorskip("PIL")
    import docx as docx_mod

    doc = docx_mod.Document()
    doc.add_paragraph("只有文字和一个小图标")
    doc.add_picture(__import__("io").BytesIO(_tiny_png()))
    f = tmp_path / "tiny.docx"
    doc.save(str(f))

    ext = get_extractor_registry().get(DOCX_MIME)
    specs = await ext.extract(f, DOCX_MIME, model=None, model_caller=None)

    assert specs[0].meta["images"] == 0
    assert "内嵌图片" not in specs[0].content


@pytest.mark.asyncio
async def test_pptx_extractor_persists_and_captions_images(tmp_path: Path):
    pytest.importorskip("pptx")
    pytest.importorskip("PIL")
    import pptx as pptx_mod

    png = _noise_png()
    prs = pptx_mod.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(0, 0, 914400, 914400)
    tb.text_frame.text = "封面标题"
    slide.shapes.add_picture(io.BytesIO(png), 0, 0)
    f = tmp_path / "deck.pptx"
    prs.save(str(f))

    stored: dict = {}

    async def fake_asset_store(filename: str, data: bytes) -> str:
        stored[filename] = data
        return f"assets/fake-{filename}"

    async def fake_caller(model, prompt, images=None):
        return "产品架构图"

    ext = get_extractor_registry().get(
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    specs = await ext.extract(
        f,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        model="vl-model",
        model_caller=fake_caller,
        asset_store=fake_asset_store,
    )

    assert len(specs) == 1
    spec = specs[0]
    assert spec.meta["images"] == 1
    assert spec.meta["slides"] == 1
    assert "封面标题" in spec.content
    assert "![图片1: slide1_" in spec.content
    assert "](assets/fake-slide1_" in spec.content
    assert "图片说明：产品架构图" in spec.content
    assert len(stored) == 1
    assert next(iter(stored.values())) == png


# ---------------------------------------------------------------------------
# ExcelExtractor: xlsx via openpyxl + legacy xls routing
# ---------------------------------------------------------------------------

XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


@pytest.mark.asyncio
async def test_excel_extractor_xlsx_renders_markdown(tmp_path: Path):
    pytest.importorskip("openpyxl")
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "负面清单"
    ws.append(["年份", "问题数|多", "来源"])
    ws.append([2023, 12, "省级检查"])
    ws.append([2024, 5, None])
    ws.append([None, None, None])
    f = tmp_path / "清单.xlsx"
    wb.save(str(f))

    ext = get_extractor_registry().get(XLSX_MIME)
    assert ext is not None and ext.name == "excel"
    specs = await ext.extract(f, XLSX_MIME, model=None, model_caller=None)

    assert len(specs) == 1
    spec = specs[0]
    assert spec.extract_mode == ExtractMode.UPLOAD
    assert spec.source_file == "清单.xlsx"
    assert spec.meta["sheets"] == 1
    assert spec.meta["rows"] == 3
    content = spec.content
    assert "## Sheet: 负面清单 (3 rows)" in content
    assert "| 年份 | 问题数\\|多 | 来源 |" in content
    assert "| 2023 | 12 | 省级检查 |" in content
    assert "| 2024 | 5 |  |" in content


@pytest.mark.asyncio
async def test_excel_extractor_xls_routes_to_xlrd_reader(tmp_path: Path, monkeypatch):
    """Legacy .xls must route to the xlrd reader (mocked — xlrd cannot write)."""
    from gyra_ext.knowledge.extractors.builtin import ExcelExtractor

    ext = get_extractor_registry().get("application/vnd.ms-excel")
    assert ext is not None and ext.name == "excel"

    f = tmp_path / "历年问题.xls"
    f.write_bytes(b"\xd0\xcf\x11\xe0fake-ole-bytes")

    def fake_read(self, path):
        return [("Sheet1", [["年份", "金额"], [2023, 1.5], [2024, 2.0]])]

    monkeypatch.setattr(ExcelExtractor, "_read_xls", fake_read)
    specs = await ext.extract(
        f, "application/vnd.ms-excel", model=None, model_caller=None
    )

    assert len(specs) == 1
    content = specs[0].content
    assert specs[0].meta["sheets"] == 1
    assert specs[0].meta["rows"] == 3
    assert "## Sheet: Sheet1 (3 rows)" in content
    assert "| 2023 | 1.5 |" in content
    assert "| 2024 | 2 |" in content


FIXTURE_XLS = Path(__file__).parent / "fixtures" / "sample.xls"


@pytest.mark.asyncio
async def test_excel_extractor_xls_real_file():
    """Regression: a real BIFF .xls must parse via xlrd's Cell attribute API."""
    pytest.importorskip("xlrd")
    assert FIXTURE_XLS.exists()
    ext = get_extractor_registry().get("application/vnd.ms-excel")
    specs = await ext.extract(
        FIXTURE_XLS, "application/vnd.ms-excel", model=None, model_caller=None
    )

    assert len(specs) == 1
    spec = specs[0]
    assert spec.meta["sheets"] == 2
    content = spec.content
    assert "## Sheet: 清单 (3 rows)" in content
    assert "| 2023 | 1.5 | 2023-06-15 00:00:00 |" in content
    assert "| 2024 | 2 | 2024-03-08 00:00:00 |" in content
    assert "## Sheet: 附表 (1 rows)" in content
    assert "| 备注 | 1 |" in content


@pytest.mark.asyncio
async def test_excel_extractor_missing_dep_raises_clear_error(
    tmp_path: Path, monkeypatch
):
    """If openpyxl isn't installed, xlsx extraction raises a clear message."""
    ext = get_extractor_registry().get(XLSX_MIME)
    f = tmp_path / "doc.xlsx"
    f.write_bytes(b"PK fake zip")

    import sys

    monkeypatch.setitem(sys.modules, "openpyxl", None)

    with pytest.raises(RuntimeError, match="openpyxl"):
        await ext.extract(f, XLSX_MIME, model=None, model_caller=None)
