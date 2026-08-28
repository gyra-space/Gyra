"""Built-in file extractors (RFC 004 §6).

Each extractor is registered via the `@extractor` decorator with one or more
mime glob patterns. Plain-text extractors (txt/md/pdf/docx/pptx) need no
model; image/audio extractors call the `model_caller` supplied by the
orchestrator.

Heavy deps (pdfplumber, python-docx, python-pptx) are imported lazily inside
the extract methods so the module imports cleanly even when optional deps
are missing — the extractor just raises a clear error at extract time.
"""

from __future__ import annotations

import asyncio
import base64
import io
import logging
import mimetypes
import os
import re
import shutil
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from gyra.knowledge.types import ExtractMode

from . import AssetStore, Extractor, ModelCaller, VerbatimSpec, extractor

logger = logging.getLogger(__name__)


def _file_mtime_iso(path: Path) -> tuple[Optional[str], Optional[int]]:
    try:
        st = path.stat()
        return (
            datetime.fromtimestamp(st.st_mtime, tz=timezone.utc).isoformat(),
            int(st.st_mtime),
        )
    except OSError:
        return None, None


# ---------------------------------------------------------------------------
# Embedded-image helpers (docx/pdf/pptx)
# ---------------------------------------------------------------------------

# Skip obvious decorations (tiny icons, spacers) — not worth an LLM call.
MIN_EMBEDDED_IMAGE_BYTES = 1024

_EMBEDDED_IMAGE_PROMPT = (
    "请用一段简明中文描述这张图片的内容：可见的文字（OCR）、图表数据、"
    "主体对象与场景。输出纯文本，不要 markdown 标题，不要解释性文字。"
)

_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"}


def _safe_image_suffix(filename: str) -> str:
    suffix = Path(filename).suffix.lower()
    return suffix if suffix in _IMAGE_SUFFIXES else ".png"


async def _store_embedded_image(
    asset_store: Optional[AssetStore], filename: str, data: bytes
) -> str:
    """Persist image bytes via the vault asset store; "" on any failure."""
    if asset_store is None or not data:
        return ""
    try:
        return await asset_store(filename, data)
    except Exception as e:  # noqa: BLE001
        logger.warning("asset_store failed for %s: %s", filename, e)
        return ""


async def _caption_embedded_image(
    filename: str,
    data: bytes,
    model: Optional[str],
    model_caller: Optional[ModelCaller],
) -> str:
    """Best-effort caption for an embedded image.

    Returns "" when no multimodal model is configured or the call fails —
    the caller keeps the image placeholder either way (never fails the
    document, per space-config decision).
    """
    if not model or not model_caller or not data:
        return ""
    try:
        with tempfile.TemporaryDirectory(prefix="ks_embed_img_") as tmp:
            tmp_path = Path(tmp) / f"img{_safe_image_suffix(filename)}"
            tmp_path.write_bytes(data)
            caption = await model_caller(model, _EMBEDDED_IMAGE_PROMPT, images=[tmp_path])
            return (caption or "").strip()
    except Exception as e:  # noqa: BLE001
        logger.warning("Embedded image caption failed for %s: %s", filename, e)
        return ""


def _image_placeholder(index: int, filename: str, ref: str, caption: str) -> str:
    """Markdown block for one embedded image."""
    label = f"图片{index}: {filename}"
    if ref:
        lines = [f"![{label}]({ref})"]
    else:
        lines = [f"（内嵌图片: {filename}，未保存）"]
    if caption:
        lines.append(f"图片说明：{caption}")
    return "\n".join(lines)


async def _process_embedded_image(
    index: int,
    filename: str,
    data: bytes,
    model: Optional[str],
    model_caller: Optional[ModelCaller],
    asset_store: Optional[AssetStore],
) -> Tuple[str, str]:
    """Store + caption + render one embedded image.

    Returns (markdown_block, caption) so callers can cache the caption for
    repeated occurrences of the same image.
    """
    ref = await _store_embedded_image(asset_store, filename, data)
    caption = await _caption_embedded_image(filename, data, model, model_caller)
    return _image_placeholder(index, filename, ref, caption), caption


# ---------------------------------------------------------------------------
# Text-family extractors
# ---------------------------------------------------------------------------


@extractor("text", ["text/*", "application/markdown", "application/json", "application/x-yaml", "application/xml"])
class TextExtractor(Extractor):
    """Plain text / markdown / json / yaml / xml — read bytes as UTF-8."""

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        raw = path.read_bytes()
        try:
            content = raw.decode("utf-8")
        except UnicodeDecodeError:
            content = raw.decode("utf-8", errors="replace")
        date_iso, mtime = _file_mtime_iso(path)
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "bytes": len(raw)},
            )
        ]


@extractor("pdf", ["application/pdf"])
class PDFExtractor(Extractor):
    """PDF → text via pdfplumber (page-separated) + embedded-image placeholders.

    Images are extracted with pypdf (PIL-backed), persisted via `asset_store`
    and captioned via `model_caller`, appended at the end of each page block.
    If pypdf/Pillow is unavailable the extractor degrades to text-only.
    """

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
        asset_store: Optional[AssetStore] = None,
    ) -> List[VerbatimSpec]:
        try:
            import pdfplumber  # type: ignore
        except ImportError as e:
            raise RuntimeError(
                "PDF extraction requires pdfplumber. Install with: pip install pdfplumber"
            ) from e

        page_images = self._extract_page_images(path)

        pages_text: List[str] = []
        img_idx = 0
        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                txt = page.extract_text() or ""
                block = f"## Page {i}\n{txt}".strip()
                for filename, data in page_images.get(i, []):
                    img_idx += 1
                    image_block, _ = await _process_embedded_image(
                        img_idx, filename, data, model, model_caller, asset_store
                    )
                    block = f"{block}\n\n{image_block}"
                pages_text.append(block)

        content = "\n\n".join(pages_text)
        date_iso, mtime = _file_mtime_iso(path)
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "pages": len(pages_text), "images": img_idx},
            )
        ]

    def _extract_page_images(self, path: Path) -> Dict[int, List[Tuple[str, bytes]]]:
        """pypdf → {page_number: [(filename, png_bytes), ...]}; best-effort."""
        try:
            from pypdf import PdfReader  # type: ignore
        except ImportError:
            logger.warning(
                "PDFExtractor: pypdf not installed, skipping embedded images in %s",
                path.name,
            )
            return {}

        out: Dict[int, List[Tuple[str, bytes]]] = {}
        try:
            reader = PdfReader(str(path))
            for pageno, page in enumerate(reader.pages, start=1):
                try:
                    for img in page.images:
                        pil = getattr(img, "image", None)
                        if pil is None:
                            continue
                        buf = io.BytesIO()
                        pil.save(buf, format="PNG")
                        data = buf.getvalue()
                        if len(data) < MIN_EMBEDDED_IMAGE_BYTES:
                            continue
                        name = getattr(img, "name", None) or f"page{pageno}_img.png"
                        out.setdefault(pageno, []).append((str(name), data))
                except Exception:  # noqa: BLE001
                    logger.debug(
                        "PDFExtractor: image extraction failed on page %d of %s",
                        pageno,
                        path.name,
                        exc_info=True,
                    )
        except Exception as e:  # noqa: BLE001
            logger.warning("PDFExtractor: pypdf failed for %s: %s", path.name, e)
        return out


@extractor("docx", [
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/msword",
])
class DocxExtractor(Extractor):
    """DOCX/DOC → text via python-docx, plus embedded-image placeholders.

    Walks every w:p element in body order (paragraphs AND table cells). For
    each paragraph, inline images are discovered via a:blip r:embed →
    image part blob, persisted through `asset_store` and captioned via
    `model_caller`; the markdown placeholder is inserted right after the
    paragraph text.
    """

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
        asset_store: Optional[AssetStore] = None,
    ) -> List[VerbatimSpec]:
        try:
            import docx  # type: ignore
        except ImportError as e:
            raise RuntimeError(
                "DOCX extraction requires python-docx. Install with: pip install python-docx"
            ) from e
        from docx.oxml.ns import qn  # type: ignore

        doc = docx.Document(str(path))
        related_parts = doc.part.related_parts
        caption_cache: Dict[str, str] = {}
        lines: List[str] = []
        img_idx = 0

        for p_el in doc.element.body.iter(qn("w:p")):
            text = "".join(t.text or "" for t in p_el.iter(qn("w:t"))).strip()
            if text:
                lines.append(text)
            for blip in p_el.iter(qn("a:blip")):
                rid = blip.get(qn("r:embed"))
                if not rid or rid not in related_parts:
                    continue
                try:
                    part = related_parts[rid]
                    data = part.blob
                    filename = str(part.partname).rsplit("/", 1)[-1]
                except Exception:  # noqa: BLE001
                    continue
                if not data or len(data) < MIN_EMBEDDED_IMAGE_BYTES:
                    continue
                img_idx += 1
                if rid in caption_cache:
                    ref = await _store_embedded_image(asset_store, filename, data)
                    block = _image_placeholder(img_idx, filename, ref, caption_cache[rid])
                else:
                    block, caption_cache[rid] = await _process_embedded_image(
                        img_idx, filename, data, model, model_caller, asset_store
                    )
                lines.append(block)

        content = "\n\n".join(lines)
        date_iso, mtime = _file_mtime_iso(path)
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "paragraphs": len(lines), "images": img_idx},
            )
        ]


@extractor("pptx", [
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
])
class PptxExtractor(Extractor):
    """PPTX/PPT → text via python-pptx (slide-separated) + image placeholders.

    Pictures are discovered recursively (including inside group shapes),
    persisted via `asset_store` and captioned via `model_caller`, appended
    after each slide's text.
    """

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
        asset_store: Optional[AssetStore] = None,
    ) -> List[VerbatimSpec]:
        try:
            import pptx  # type: ignore
        except ImportError as e:
            raise RuntimeError(
                "PPTX extraction requires python-pptx. Install with: pip install python-pptx"
            ) from e
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore

        def iter_pictures(shapes):
            for shape in shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                        yield from iter_pictures(shape.shapes)
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        yield shape
                except Exception:  # noqa: BLE001
                    continue

        prs = pptx.Presentation(str(path))
        slides_text: List[str] = []
        img_idx = 0
        for i, slide in enumerate(prs.slides, start=1):
            texts: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            texts.append(t)
            block = f"## Slide {i}\n" + "\n".join(texts) if texts else f"## Slide {i}"
            for pic in iter_pictures(slide.shapes):
                try:
                    image = pic.image
                    data = image.blob
                    ext = (image.ext or "png").lower()
                    if ext not in _IMAGE_SUFFIXES:
                        ext = "png"
                    filename = f"slide{i}_{pic.shape_id}.{ext}"
                except Exception:  # noqa: BLE001
                    continue
                if not data or len(data) < MIN_EMBEDDED_IMAGE_BYTES:
                    continue
                img_idx += 1
                image_block, _ = await _process_embedded_image(
                    img_idx, filename, data, model, model_caller, asset_store
                )
                block = f"{block}\n\n{image_block}"
            slides_text.append(block)

        content = "\n\n".join(slides_text)
        date_iso, mtime = _file_mtime_iso(path)
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "slides": len(slides_text), "images": img_idx},
            )
        ]


# ---------------------------------------------------------------------------
# Multimodal extractors (require model_caller)
# ---------------------------------------------------------------------------


@extractor("image", ["image/*"])
class ImageExtractor(Extractor):
    """Image → caption verbatim via a multimodal LLM.

    Calls `model_caller(model, prompt, images=[path])`. Falls back to an
    OCR-only verbatim (empty caption, marked deprecated) if no model_caller.
    """

    PROMPT = (
        "请详细描述这张图片的内容。包括：可见的文字（OCR）、主体对象、场景、"
        "图表数据、以及任何对知识库有用的元信息。输出纯文本，不要 markdown 标题。"
    )

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        date_iso, mtime = _file_mtime_iso(path)
        if not model or not model_caller:
            logger.warning(
                "ImageExtractor: no model/model_caller supplied for %s; "
                "writing empty-caption verbatim",
                path.name,
            )
            caption = ""
        else:
            try:
                caption = await model_caller(model, self.PROMPT, images=[path])
            except Exception as e:
                logger.exception("ImageExtractor model call failed for %s", path.name)
                raise RuntimeError(f"Image caption model call failed: {e}") from e

        content = (
            f"[Image: {path.name}]\n\n"
            f"**MIME**: {mime}\n\n"
            f"**Caption**:\n{caption}"
        )
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "model": model, "caption_len": len(caption)},
            )
        ]


@extractor("audio", ["audio/*"])
class AudioExtractor(Extractor):
    """Audio → transcription verbatim via an ASR/multimodal model.

    Calls `model_caller(model, prompt, images=[path])` — the caller is
    responsible for routing audio files to a model that accepts them
    (e.g., qwen-audio). If no model_caller is supplied, raises a clear error.
    """

    PROMPT = "请转录这段音频的内容。保留原文，不要总结。如果有多人对话，标注说话人。"

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        date_iso, mtime = _file_mtime_iso(path)
        if not model or not model_caller:
            raise RuntimeError(
                f"Audio extraction for {path.name} requires a multimodal/ASR model. "
                "Configure a multimodal_model on the space or pass a model_override."
            )
        try:
            transcript = await model_caller(model, self.PROMPT, images=[path])
        except Exception as e:
            logger.exception("AudioExtractor model call failed for %s", path.name)
            raise RuntimeError(f"Audio transcription model call failed: {e}") from e

        content = (
            f"[Audio: {path.name}]\n\n"
            f"**MIME**: {mime}\n\n"
            f"**Transcript**:\n{transcript}"
        )
        return [
            VerbatimSpec(
                content=content,
                source_file=path.name,
                extract_mode=ExtractMode.UPLOAD,
                content_date=date_iso,
                source_mtime=mtime,
                meta={"mime": mime, "model": model},
            )
        ]


# ---------------------------------------------------------------------------
# Video
# ---------------------------------------------------------------------------


async def _run_cmd(cmd: List[str]) -> int:
    """Run a subprocess quietly; returns the exit code."""
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.DEVNULL,
        stderr=asyncio.subprocess.DEVNULL,
    )
    return await proc.wait()


@extractor("video", ["video/*"])
class VideoExtractor(Extractor):
    """Video → 双策略:原生视频模型直连,不支持/失败则降级 ffmpeg 抽帧管线。

    - 原生直连:model 命中 NATIVE_MODEL_PATTERN(gemini / qwen-vl / qwen-omni)
      且文件 ≤ MAX_NATIVE_BYTES 时,经 model_caller(videos=[path]) 直接理解;
    - 抽帧管线:ffmpeg 按间隔抽帧(cap MAX_FRAMES)→ 分批走 image 通道逐批
      描述;音轨抽出转写(best-effort);合并为带片段结构的 verbatim。

    环境依赖:抽帧管线需要 ffmpeg 在 PATH 上。
    """

    NATIVE_MODEL_PATTERN = re.compile(r"gemini|qwen.*(vl|omni)", re.IGNORECASE)
    MAX_NATIVE_BYTES = 20 * 1024 * 1024
    MAX_FRAMES = 24
    FRAME_BATCH = 8
    FRAME_INTERVAL_SEC = 5

    PROMPT_NATIVE = (
        "请详细描述这个视频的内容:场景、事件发展过程、可见的文字、"
        "以及任何对知识库有用的关键信息。输出纯文本,不要 markdown 标题。"
    )
    PROMPT_FRAMES = (
        "这是同一个视频按时间顺序抽取的一组帧。请描述这组帧里发生的事情、"
        "场景变化、可见的文字与关键信息。输出纯文本,不要 markdown 标题。"
    )
    PROMPT_AUDIO = "请转录这段音频的内容。保留原文,不要总结。如果有多人对话,标注说话人。"

    async def extract(
        self,
        path: Path,
        mime: str,
        model: Optional[str],
        model_caller: Optional[ModelCaller],
    ) -> List[VerbatimSpec]:
        if not model or not model_caller:
            raise RuntimeError(
                f"Video extraction for {path.name} requires a multimodal model. "
                "Configure a multimodal_model on the space or pass a model_override."
            )
        date_iso, mtime = _file_mtime_iso(path)

        # 1. 原生视频模型直连(限小文件;失败/空输出降级抽帧)
        try:
            size = path.stat().st_size
        except OSError:
            size = 0
        if self.NATIVE_MODEL_PATTERN.search(model) and size <= self.MAX_NATIVE_BYTES:
            try:
                text = await model_caller(model, self.PROMPT_NATIVE, videos=[path])
                if text and text.strip():
                    content = (
                        f"[Video: {path.name}]\n\n"
                        f"**MIME**: {mime}\n\n"
                        f"**Content**:\n{text}"
                    )
                    return [
                        VerbatimSpec(
                            content=content,
                            source_file=path.name,
                            extract_mode=ExtractMode.UPLOAD,
                            content_date=date_iso,
                            source_mtime=mtime,
                            meta={"mime": mime, "model": model, "strategy": "native"},
                        )
                    ]
                logger.warning(
                    "VideoExtractor: native call returned empty for %s; "
                    "falling back to frame pipeline",
                    path.name,
                )
            except Exception as e:  # noqa: BLE001
                logger.warning(
                    "VideoExtractor: native call failed for %s (%s); "
                    "falling back to frame pipeline",
                    path.name,
                    e,
                )

        # 2. ffmpeg 抽帧 + 音轨管线
        if not shutil.which("ffmpeg"):
            raise RuntimeError(
                "Video extraction requires ffmpeg (not found on PATH), "
                "or a native video model (gemini / qwen-vl)."
            )
        import tempfile

        with tempfile.TemporaryDirectory(prefix="ks_video_") as tmp:
            tmp_dir = Path(tmp)
            frames = await self._extract_frames(path, tmp_dir)
            segments: List[str] = []
            for i in range(0, len(frames), self.FRAME_BATCH):
                batch = frames[i : i + self.FRAME_BATCH]
                desc = await model_caller(model, self.PROMPT_FRAMES, images=batch)
                segments.append(f"[片段 {i // self.FRAME_BATCH + 1}]\n{desc}")

            transcript = ""
            audio_path = await self._extract_audio(path, tmp_dir)
            if audio_path is not None:
                try:
                    transcript = await model_caller(
                        model, self.PROMPT_AUDIO, images=[audio_path]
                    )
                except Exception as e:  # noqa: BLE001
                    logger.warning(
                        "VideoExtractor: audio transcription failed for %s: %s",
                        path.name,
                        e,
                    )

            if not segments and not transcript.strip():
                raise RuntimeError(
                    f"Video extraction for {path.name} produced no content "
                    "(no frames, no audio track)"
                )
            parts = [f"[Video: {path.name}]\n\n**MIME**: {mime}"]
            if segments:
                parts.append("**画面**:\n" + "\n\n".join(segments))
            if transcript.strip():
                parts.append("**音轨转写**:\n" + transcript)
            return [
                VerbatimSpec(
                    content="\n\n".join(parts),
                    source_file=path.name,
                    extract_mode=ExtractMode.UPLOAD,
                    content_date=date_iso,
                    source_mtime=mtime,
                    meta={
                        "mime": mime,
                        "model": model,
                        "strategy": "frames",
                        "frames": len(frames),
                        "audio": bool(transcript.strip()),
                    },
                )
            ]

    async def _extract_frames(self, path: Path, out_dir: Path) -> List[Path]:
        """ffmpeg 按固定间隔抽帧(缩放至宽 768),上限 MAX_FRAMES。"""
        pattern = out_dir / "frame_%03d.jpg"
        rc = await _run_cmd(
            [
                "ffmpeg",
                "-y",
                "-i",
                str(path),
                "-vf",
                f"fps=1/{self.FRAME_INTERVAL_SEC},scale=768:-2",
                "-frames:v",
                str(self.MAX_FRAMES),
                str(pattern),
            ]
        )
        if rc != 0:
            logger.warning("VideoExtractor: ffmpeg frame extraction failed for %s", path.name)
            return []
        return sorted(out_dir.glob("frame_*.jpg"))

    async def _extract_audio(self, path: Path, out_dir: Path) -> Optional[Path]:
        """ffmpeg 抽出音轨为 16k 单声道 wav;无音轨返回 None。"""
        out = out_dir / "audio.wav"
        rc = await _run_cmd(
            ["ffmpeg", "-y", "-i", str(path), "-vn", "-ac", "1", "-ar", "16000", str(out)]
        )
        if rc != 0 or not out.exists() or out.stat().st_size == 0:
            return None
        return out
